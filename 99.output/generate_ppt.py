from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI 기반 모의 면접 시스템"
    subtitle.text = "상세 업무 흐름도 (Workflow) - Phase별 상세"

    # Helper function to add a content slide
    def add_phase_slide(phase_title, steps):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = phase_title
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for step in steps:
            p = tf.add_paragraph()
            p.text = step
            p.font.size = Pt(20)
            p.level = 0

    # Phase 1
    add_phase_slide(
        "Phase 1: 인풋 및 컨텍스트 설정",
        [
            "1. User Input (시작)",
            "   - 사용자 자기소개서 및 지원 직무 정보 입력",
            "",
            "2. Job Analysis",
            "   - 채용 공고/직무에서 핵심 역량 키워드 추출",
            "",
            "3. Question Generation (LLM)",
            "   - 맞춤형 면접 질문 생성",
            "   - 질문 큐(Queue)에 적재"
        ]
    )

    # Phase 2
    add_phase_slide(
        "Phase 2: 실시간 면접 엔진",
        [
            "1. Session Start",
            "   - 큐에서 질문을 가져와 사용자에게 제시",
            "",
            "2. Media Streaming",
            "   - WebRTC를 통한 실시간 오디오/비디오 수신",
            "",
            "3. Multi-modal Analysis",
            "   - STT: 음성을 텍스트로 변환",
            "   - Vision: 표정 및 시선 처리 분석",
            "   - Audio: 목소리 톤 및 속도 분석",
            "",
            "4. Intent Analysis",
            "   - 답변 완료 여부 및 주제 이탈 감지"
        ]
    )

    # Phase 3
    add_phase_slide(
        "Phase 3: 동적 질문 제어",
        [
            "1. Tail Check (꼬리 질문 판단)",
            "   - 답변에서 특정 키워드나 모호함이 감지되면 분기",
            "",
            "2. Tail-Question Logic",
            "   - 심층 질문(Deep Dive) 생성하여 큐에 추가",
            "",
            "3. Next / End Loop",
            "   - 다음 질문이 있는지 확인",
            "   - 없으면 면접 종료(End Session)"
        ]
    )

    # Phase 4
    add_phase_slide(
        "Phase 4: 평가 및 리포팅",
        [
            "1. Score Aggregation",
            "   - 답변 내용 점수 (70%) + 태도 점수 (30%) 합산",
            "",
            "2. Visualization",
            "   - 평가 결과 시각화 (방사형 차트 등)",
            "   - 영상 타임스탬프 매핑",
            "",
            "3. Final Report",
            "   - 최종 피드백 리포트 생성 및 사용자 제공"
        ]
    )

    output_path = "Workflow_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
